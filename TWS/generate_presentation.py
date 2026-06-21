"""
The Wolf Stack — FYP Viva Presentation Generator
Generates TheWolfStack_FYP_Presentation.pptx
Run: python generate_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
BG         = RGBColor(0x0F, 0x17, 0x2A)   # dark navy
ACCENT     = RGBColor(0x3B, 0x82, 0xF6)   # blue
ORANGE     = RGBColor(0xF9, 0x73, 0x16)   # orange highlight
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCB, 0xD5, 0xE1)
CARD_BG    = RGBColor(0x1E, 0x29, 0x3B)   # slightly lighter navy for cards
GREEN      = RGBColor(0x22, 0xC5, 0x5E)
RED        = RGBColor(0xEF, 0x44, 0x44)

W = Inches(13.33)
H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]          # completely blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, l, t, w, h,
                font_size=Pt(18), bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, word_wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def add_accent_bar(slide, color=ACCENT):
    """Thin horizontal accent line near top."""
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill=color)


def slide_heading(slide, text, top=Inches(0.35), color=WHITE, size=Pt(32)):
    add_textbox(slide, text,
                Inches(0.5), top, Inches(12.3), Inches(0.7),
                font_size=size, bold=True, color=color)


def bullet_block(slide, items, top, left=Inches(0.6), width=Inches(12.2),
                 size=Pt(19), bullet="•", color=WHITE, spacing=Inches(0.42)):
    for i, item in enumerate(items):
        text = f"{bullet}  {item}"
        add_textbox(slide, text,
                    left, top + i * spacing, width, Inches(0.5),
                    font_size=size, color=color)


# ── Individual slide builders ─────────────────────────────────────────────────

def slide_01_title(prs):
    s = blank_slide(prs)
    fill_bg(s)

    # top gradient bar
    add_rect(s, Inches(0), Inches(0), W, Inches(0.08), fill=ACCENT)
    add_rect(s, Inches(0), Inches(0.08), W, Inches(0.04), fill=ORANGE)

    # decorative right panel
    add_rect(s, Inches(9.8), Inches(0.12), Inches(3.53), H - Inches(0.12), fill=CARD_BG)

    # big title
    add_textbox(s, "The Wolf Stack",
                Inches(0.6), Inches(1.4), Inches(9.0), Inches(1.4),
                font_size=Pt(62), bold=True, color=WHITE)

    # subtitle
    add_textbox(s, "Multi-Tenant SaaS ERP Portal",
                Inches(0.6), Inches(2.85), Inches(9.0), Inches(0.7),
                font_size=Pt(28), bold=False, color=ACCENT)

    # divider line
    add_rect(s, Inches(0.6), Inches(3.65), Inches(5.0), Inches(0.04), fill=ORANGE)

    # tag line
    add_textbox(s, "Final Year Project  ·  Bachelor of Information Technology",
                Inches(0.6), Inches(3.8), Inches(9.0), Inches(0.5),
                font_size=Pt(16), color=LIGHT_GREY)

    # right panel content
    for i, line in enumerate([
        "Stack",
        "React 18  +  Node.js",
        "MongoDB  +  Express",
        "JWT  ·  RBAC  ·  AWS S3",
        "Socket.io  ·  Railway",
    ]):
        clr = ORANGE if i == 0 else LIGHT_GREY
        sz  = Pt(18) if i == 0 else Pt(15)
        bd  = (i == 0)
        add_textbox(s, line, Inches(10.0), Inches(1.5 + i * 0.7), Inches(3.1), Inches(0.5),
                    font_size=sz, bold=bd, color=clr)

    # bottom team block
    add_rect(s, Inches(0), H - Inches(1.3), W, Inches(1.3), fill=CARD_BG)
    add_textbox(s, "[Your Name]  |  [Member 2]  |  [Member 3]",
                Inches(0.6), H - Inches(1.2), Inches(8), Inches(0.45),
                font_size=Pt(16), bold=True, color=WHITE)
    add_textbox(s, "Supervised by: [Supervisor Name]     ·     [University Name] — Dept. of Information Technology     ·     2025–2026",
                Inches(0.6), H - Inches(0.75), Inches(12), Inches(0.4),
                font_size=Pt(13), color=LIGHT_GREY)


def slide_02_problem(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_accent_bar(s, RED)
    slide_heading(s, "The Problem", color=WHITE)

    add_textbox(s, "Why do software houses struggle with operations management?",
                Inches(0.6), Inches(1.1), Inches(12), Inches(0.45),
                font_size=Pt(16), color=LIGHT_GREY, italic=True)

    problems = [
        ("❌  Scattered Tools",
         "HR in Excel, Projects on WhatsApp, Finance in manual registers — no single system"),
        ("❌  Unaffordable Enterprise ERP",
         "SAP, Oracle, Microsoft Dynamics cost thousands of dollars — far beyond SME budgets"),
        ("❌  No Local SaaS ERP",
         "No affordable Pakistani multi-tenant SaaS ERP built for software houses exists"),
        ("❌  High Setup Cost per Company",
         "Every business needs separate deployments — no shared infrastructure, no economies of scale"),
    ]

    for i, (heading, detail) in enumerate(problems):
        top = Inches(1.65 + i * 1.35)
        add_rect(s, Inches(0.5), top, Inches(12.33), Inches(1.15), fill=CARD_BG)
        add_textbox(s, heading, Inches(0.75), top + Inches(0.1), Inches(11.5), Inches(0.38),
                    font_size=Pt(18), bold=True, color=RED)
        add_textbox(s, detail,  Inches(0.75), top + Inches(0.48), Inches(11.5), Inches(0.55),
                    font_size=Pt(15), color=LIGHT_GREY)


def slide_03_solution(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_accent_bar(s, GREEN)
    slide_heading(s, "Our Solution — Introducing The Wolf Stack", color=WHITE)

    add_textbox(s, "One platform. Every company. Fully isolated.",
                Inches(0.6), Inches(1.1), Inches(12), Inches(0.45),
                font_size=Pt(18), bold=True, color=ACCENT)

    solutions = [
        ("✅  Multi-Tenant SaaS Platform",
         "One deployed instance serves many companies — each with fully isolated data & configuration"),
        ("✅  Complete ERP in One Place",
         "HR & Attendance · Finance & Payroll · Project Management — all under one roof"),
        ("✅  Two-Level Administration",
         "Supra Admin manages the entire platform; each Tenant Admin manages their own workspace"),
        ("✅  Role-Based Access Control",
         "Owner → Admin → Manager → HR → Finance → Employee — granular permissions at every level"),
        ("✅  Cloud-Native & Affordable",
         "Deployed on Railway (cloud), MongoDB Atlas — no expensive infrastructure for clients"),
    ]

    for i, (heading, detail) in enumerate(solutions):
        top = Inches(1.65 + i * 1.12)
        add_rect(s, Inches(0.5), top, Inches(12.33), Inches(0.98), fill=CARD_BG)
        add_textbox(s, heading, Inches(0.75), top + Inches(0.08), Inches(11.5), Inches(0.35),
                    font_size=Pt(17), bold=True, color=GREEN)
        add_textbox(s, detail,  Inches(0.75), top + Inches(0.44), Inches(11.5), Inches(0.45),
                    font_size=Pt(14), color=LIGHT_GREY)


def slide_04_architecture(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_accent_bar(s, ACCENT)
    slide_heading(s, "System Architecture")

    # Three column cards
    cols = [
        ("Frontend", ACCENT, [
            "React 18  +  Vite",
            "Tailwind CSS  +  Ant Design",
            "Context API (Auth, Socket)",
            "Axios  +  React Router v6",
            "Deployed on Railway",
        ]),
        ("Backend", ORANGE, [
            "Node.js  +  Express.js",
            "6 Modular Route Groups",
            "JWT  +  RBAC Middleware",
            "Socket.io (real-time)",
            "Deployed on Railway",
        ]),
        ("Database", RGBColor(0xA7, 0x8B, 0xFA), [
            "MongoDB Atlas (cloud)",
            "Mongoose ODM v7",
            "88+ Schema Models",
            "orgId-scoped queries",
            "AWS S3 (file storage)",
        ]),
    ]

    col_w = Inches(3.9)
    for ci, (title, color, bullets) in enumerate(cols):
        lft = Inches(0.4 + ci * 4.3)
        add_rect(s, lft, Inches(1.15), col_w, Inches(5.85), fill=CARD_BG)
        # colour top strip
        add_rect(s, lft, Inches(1.15), col_w, Inches(0.08), fill=color)
        add_textbox(s, title, lft + Inches(0.15), Inches(1.28), col_w - Inches(0.3), Inches(0.5),
                    font_size=Pt(20), bold=True, color=color)
        for bi, b in enumerate(bullets):
            add_textbox(s, f"→  {b}",
                        lft + Inches(0.15),
                        Inches(1.95 + bi * 0.85),
                        col_w - Inches(0.3),
                        Inches(0.55),
                        font_size=Pt(15), color=WHITE)

    # flow arrow label at bottom
    add_textbox(s, "Request Flow:   Browser  →  HTTPS  →  Express Middleware (Auth · RBAC · Tenant Filter)  →  MongoDB",
                Inches(0.5), Inches(7.05), Inches(12.5), Inches(0.35),
                font_size=Pt(13), color=LIGHT_GREY, italic=True)


def slide_05_techstack(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_accent_bar(s, ACCENT)
    slide_heading(s, "Technology Stack")

    rows = [
        ("Frontend",    "React 18,  Tailwind CSS,  Ant Design,  Mantine UI,  Framer Motion"),
        ("Backend",     "Node.js,  Express.js  (modular architecture)"),
        ("Database",    "MongoDB Atlas,  Mongoose ODM 7"),
        ("Auth",        "JWT (HttpOnly Cookies),  bcrypt,  Passport.js"),
        ("Storage",     "AWS S3  (file & image uploads via Multer)"),
        ("Real-time",   "Socket.io 4  (live notifications & collaboration)"),
        ("Security",    "Helmet,  Rate Limiting,  RBAC / PBAC,  Token Blacklist"),
        ("Monitoring",  "Prometheus,  Winston Logging,  Sentry Error Tracking"),
        ("Deployment",  "Railway  (Frontend + Backend),  MongoDB Atlas  (Database)"),
    ]

    col_colors = [ACCENT, ORANGE]
    for i, (layer, tech) in enumerate(rows):
        top = Inches(1.2 + i * 0.68)
        bg  = CARD_BG if i % 2 == 0 else RGBColor(0x16, 0x21, 0x33)
        add_rect(s, Inches(0.4), top, Inches(12.5), Inches(0.62), fill=bg)
        add_textbox(s, layer, Inches(0.55), top + Inches(0.1), Inches(2.2), Inches(0.42),
                    font_size=Pt(15), bold=True, color=ACCENT)
        add_textbox(s, tech,  Inches(2.85), top + Inches(0.1), Inches(10.0), Inches(0.42),
                    font_size=Pt(15), color=WHITE)


def slide_06_multitenancy(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_accent_bar(s, ORANGE)
    slide_heading(s, "Multi-Tenancy — The Core Innovation")

    add_textbox(s, "One platform instance  ·  Many isolated companies  ·  Zero data leakage",
                Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
                font_size=Pt(16), bold=True, color=ORANGE)

    layers = [
        ("1  ·  Query-Level Isolation",   ACCENT,  "orgId auto-injected into every MongoDB query via middleware — no manual filtering needed"),
        ("2  ·  Token-Level Isolation",   GREEN,   "JWT claims verified against DB; tenant slug in URL validated on every request"),
        ("3  ·  Access-Level Isolation",  ORANGE,  "RBAC ensures users can only read/write data belonging to their own tenant"),
        ("4  ·  Socket-Level Isolation",  RGBColor(0xA7,0x8B,0xFA), "Real-time events scoped to tenant:{tenantId} rooms — no cross-tenant leakage"),
    ]

    for i, (label, color, desc) in enumerate(layers):
        top = Inches(1.65 + i * 1.12)
        add_rect(s, Inches(0.5), top, Inches(12.33), Inches(0.98), fill=CARD_BG)
        add_rect(s, Inches(0.5), top, Inches(0.06), Inches(0.98), fill=color)
        add_textbox(s, label, Inches(0.72), top + Inches(0.08), Inches(4.0), Inches(0.38),
                    font_size=Pt(16), bold=True, color=color)
        add_textbox(s, desc,  Inches(0.72), top + Inches(0.48), Inches(11.5), Inches(0.42),
                    font_size=Pt(14), color=LIGHT_GREY)

    add_textbox(s, "Tenant Types:   Software House  ·  Business  ·  Warehouse",
                Inches(0.6), Inches(7.05), Inches(12), Inches(0.35),
                font_size=Pt(14), color=LIGHT_GREY, italic=True)


def _module_slide(prs, title, color, icon_label, features_left, features_right):
    s = blank_slide(prs)
    fill_bg(s)
    add_accent_bar(s, color)
    slide_heading(s, title)

    # module badge
    add_rect(s, Inches(10.5), Inches(0.22), Inches(2.6), Inches(0.5), fill=color)
    add_textbox(s, icon_label, Inches(10.5), Inches(0.22), Inches(2.6), Inches(0.5),
                font_size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    mid = Inches(7.0)
    for i, feat in enumerate(features_left):
        top = Inches(1.25 + i * 0.83)
        add_rect(s, Inches(0.4), top, Inches(6.3), Inches(0.72), fill=CARD_BG)
        add_rect(s, Inches(0.4), top, Inches(0.06), Inches(0.72), fill=color)
        add_textbox(s, feat, Inches(0.62), top + Inches(0.16), Inches(5.9), Inches(0.4),
                    font_size=Pt(15), color=WHITE)

    for i, feat in enumerate(features_right):
        top = Inches(1.25 + i * 0.83)
        add_rect(s, Inches(6.95), top, Inches(6.3), Inches(0.72), fill=CARD_BG)
        add_rect(s, Inches(6.95), top, Inches(0.06), Inches(0.72), fill=color)
        add_textbox(s, feat, Inches(7.17), top + Inches(0.16), Inches(5.9), Inches(0.4),
                    font_size=Pt(15), color=WHITE)


def slide_07_hr(prs):
    _module_slide(prs,
        "Module  ·  HR & Attendance Management", GREEN, "HR MODULE",
        features_left=[
            "Employee Records & Profiles",
            "Attendance Tracking (shifts & policies)",
            "Leave Request & Approval Workflows",
            "Payroll Processing & Payslip Generation",
        ],
        features_right=[
            "Performance Reviews & Employee Metrics",
            "Onboarding & Training Management",
            "Recruitment Pipeline",
            "AI-Assisted Payroll Calculation",
        ]
    )


def slide_08_finance(prs):
    _module_slide(prs,
        "Module  ·  Finance & Payroll", ORANGE, "FINANCE MODULE",
        features_left=[
            "Accounts Payable & Receivable (AP/AR)",
            "Chart of Accounts",
            "Expense Tracking & Approval",
            "Invoice Generation & Billing",
        ],
        features_right=[
            "Subscription Plan Management",
            "Equity Tracking",
            "Financial Reports (Excel / PDF export)",
            "Dashboard with Key Metrics",
        ]
    )


def slide_09_projects(prs):
    _module_slide(prs,
        "Module  ·  Project Management", ACCENT, "PM MODULE",
        features_left=[
            "Project Creation with Access Control",
            "Task Management + Dependencies",
            "Kanban Board (Drag & Drop)",
            "Sprint Planning  (Agile / Scrum)",
        ],
        features_right=[
            "Gantt Chart & Timeline View",
            "Milestones & Deliverables",
            "Time Tracking & Dev Metrics",
            "Change Request Management",
        ]
    )


def slide_10_supra(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_accent_bar(s, RGBColor(0xA7, 0x8B, 0xFA))
    slide_heading(s, "Supra Admin Panel — Platform Control")

    add_textbox(s, "The master administrator who controls the entire Wolf Stack platform",
                Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
                font_size=Pt(16), color=LIGHT_GREY, italic=True)

    items = [
        ("Platform Dashboard",    "Live overview of all tenants, system health, and usage metrics"),
        ("Tenant Management",     "Approve, suspend, configure, and provision new tenant companies"),
        ("ERP Module Assignment", "Enable / disable specific ERP modules per tenant based on their plan"),
        ("Master ERP Templates",  "Pre-built department & role templates for rapid tenant setup"),
        ("Platform Admin Roles",  "Separate roles: Platform Admin, Support, Billing Analyst"),
        ("System-Wide Audit Log", "Every platform-level action logged with timestamp, IP, and user"),
    ]

    col_w = Inches(6.0)
    for i, (heading, detail) in enumerate(items):
        col = i % 2
        row = i // 2
        lft = Inches(0.4 + col * 6.55)
        top = Inches(1.65 + row * 1.72)
        add_rect(s, lft, top, col_w, Inches(1.52), fill=CARD_BG)
        add_rect(s, lft, top, col_w, Inches(0.06), fill=RGBColor(0xA7, 0x8B, 0xFA))
        add_textbox(s, heading, lft + Inches(0.15), top + Inches(0.18), col_w - Inches(0.3), Inches(0.42),
                    font_size=Pt(16), bold=True, color=RGBColor(0xA7, 0x8B, 0xFA))
        add_textbox(s, detail,  lft + Inches(0.15), top + Inches(0.65), col_w - Inches(0.3), Inches(0.75),
                    font_size=Pt(14), color=WHITE)


def slide_11_security(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_accent_bar(s, RED)
    slide_heading(s, "Security Implementation")

    items = [
        (RED,    "Authentication",     "JWT stored in HttpOnly cookies — eliminates XSS token theft from localStorage"),
        (ORANGE, "Authorization",      "RBAC + PBAC — 10+ defined roles per tenant with resource-action permission model"),
        (GREEN,  "Rate Limiting",      "100 req/15min globally; auth routes protected with strict 5 req/15min limits"),
        (ACCENT, "Audit Logging",      "Every sensitive action logged with: user ID, IP address, timestamp, action type"),
        (RGBColor(0xA7,0x8B,0xFA), "Input Validation", "express-validator on all POST/PUT routes — prevents injection attacks"),
        (ORANGE, "Token Blacklisting", "Revoked / logout tokens tracked in DB — replay attacks prevented"),
    ]

    col_w = Inches(6.0)
    for i, (color, heading, detail) in enumerate(items):
        col = i % 2
        row = i // 2
        lft = Inches(0.4 + col * 6.55)
        top = Inches(1.2 + row * 2.05)
        add_rect(s, lft, top, col_w, Inches(1.82), fill=CARD_BG)
        add_rect(s, lft, top, Inches(0.06), Inches(1.82), fill=color)
        add_textbox(s, heading, lft + Inches(0.2), top + Inches(0.15), col_w - Inches(0.35), Inches(0.42),
                    font_size=Pt(17), bold=True, color=color)
        add_textbox(s, detail,  lft + Inches(0.2), top + Inches(0.62), col_w - Inches(0.35), Inches(1.0),
                    font_size=Pt(14), color=WHITE)


def slide_12_challenges(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_accent_bar(s, ORANGE)
    slide_heading(s, "Challenges & How We Solved Them")

    challenges = [
        ("Cross-Tenant Data Leakage",
         "Auto-inject orgId query filter middleware on every API route — no manual tenant checks required"),
        ("JWT Security in Browser",
         "Switched from localStorage to HttpOnly cookies — tokens invisible to JavaScript, eliminating XSS risk"),
        ("Dynamic ERP Modules per Tenant",
         "Tenant record stores enabled module flags; frontend reads config and renders only permitted modules"),
        ("Role Complexity Across Tenants",
         "Unified RBAC with tenant-scoped role assignments — same role system, different data per tenant"),
        ("Scalable Modular Backend",
         "Express routes split into 6 independent modules — each deployable and testable in isolation"),
    ]

    for i, (challenge, solution) in enumerate(challenges):
        top = Inches(1.3 + i * 1.22)
        add_rect(s, Inches(0.4),  top, Inches(5.5), Inches(1.06), fill=CARD_BG)
        add_rect(s, Inches(5.95), top, Inches(7.0), Inches(1.06), fill=RGBColor(0x16, 0x21, 0x33))

        add_textbox(s, f"⚡  {challenge}",
                    Inches(0.55), top + Inches(0.1), Inches(5.2), Inches(0.42),
                    font_size=Pt(14), bold=True, color=ORANGE)
        add_textbox(s, f"✅  {solution}",
                    Inches(6.1), top + Inches(0.1), Inches(6.65), Inches(0.82),
                    font_size=Pt(14), color=WHITE)

    # column headers
    add_textbox(s, "CHALLENGE", Inches(0.55), Inches(1.08), Inches(3), Inches(0.3),
                font_size=Pt(12), bold=True, color=ORANGE)
    add_textbox(s, "SOLUTION", Inches(6.1), Inches(1.08), Inches(3), Inches(0.3),
                font_size=Pt(12), bold=True, color=GREEN)


def slide_13_future(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_accent_bar(s, ACCENT)
    slide_heading(s, "Future Roadmap & Thank You")

    roadmap = [
        ("📱  Mobile App",           "React Native app for employees — attendance, tasks, payslips on the go"),
        ("🤖  AI Analytics",         "AI-powered business forecasting, automated payroll insights, anomaly detection"),
        ("💬  WhatsApp / SMS",        "Notifications via WhatsApp Business API and SMS for approvals & alerts"),
        ("🔌  Plugin Marketplace",   "Tenants can install third-party ERP plugins — extensible module system"),
        ("🏥  New Tenant Types",     "Expand beyond Software Houses to Healthcare, Education, and Retail sectors"),
    ]

    for i, (title, desc) in enumerate(roadmap):
        top = Inches(1.3 + i * 1.0)
        add_rect(s, Inches(0.4), top, Inches(11.0), Inches(0.88), fill=CARD_BG)
        add_textbox(s, title, Inches(0.6), top + Inches(0.07), Inches(3.5), Inches(0.38),
                    font_size=Pt(16), bold=True, color=ACCENT)
        add_textbox(s, desc,  Inches(4.0), top + Inches(0.07), Inches(7.2), Inches(0.72),
                    font_size=Pt(15), color=WHITE)

    # Thank you footer
    add_rect(s, Inches(0), H - Inches(1.1), W, Inches(1.1), fill=RGBColor(0x1E, 0x40, 0xAF))
    add_textbox(s, "Thank You  —  Questions Welcome",
                Inches(0), H - Inches(1.0), W, Inches(0.55),
                font_size=Pt(28), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, "The Wolf Stack  ·  Multi-Tenant SaaS ERP Portal  ·  FYP 2025–2026",
                Inches(0), H - Inches(0.45), W, Inches(0.35),
                font_size=Pt(14), color=LIGHT_GREY, align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    print("Building slides...")
    slide_01_title(prs);       print("  [1/13] Title slide")
    slide_02_problem(prs);     print("  [2/13] The Problem")
    slide_03_solution(prs);    print("  [3/13] Our Solution")
    slide_04_architecture(prs);print("  [4/13] System Architecture")
    slide_05_techstack(prs);   print("  [5/13] Technology Stack")
    slide_06_multitenancy(prs);print("  [6/13] Multi-Tenancy")
    slide_07_hr(prs);          print("  [7/13] HR & Attendance Module")
    slide_08_finance(prs);     print("  [8/13] Finance & Payroll Module")
    slide_09_projects(prs);    print("  [9/13] Project Management Module")
    slide_10_supra(prs);       print(" [10/13] Supra Admin Panel")
    slide_11_security(prs);    print(" [11/13] Security Implementation")
    slide_12_challenges(prs);  print(" [12/13] Challenges & Solutions")
    slide_13_future(prs);      print(" [13/13] Future Roadmap & Thank You")

    out = "TheWolfStack_FYP_Presentation.pptx"
    prs.save(out)
    print(f"\nDONE  ->  {out}")
    print("   Open in PowerPoint and fill in [Your Name], [Supervisor], [University].")


if __name__ == "__main__":
    main()
